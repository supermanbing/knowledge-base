from __future__ import annotations

import csv
import html.parser
import io
import json
import re
from pathlib import Path
from typing import Any

try:
    import fitz
except ImportError:
    fitz = None

try:
    from docx import Document
    from docx.document import Document as DocumentType
    from docx.table import Table
    from docx.text.paragraph import Paragraph
except ImportError:
    Document = None
    DocumentType = Any
    Table = Any
    Paragraph = Any

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    RecursiveCharacterTextSplitter = None

from backend.multimodal import describe_image

# ===== 支持的格式 =====
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".vue",
    ".java", ".cpp", ".c", ".h", ".hpp", ".rs", ".go",
    ".rb", ".php", ".sh", ".bash", ".zsh",
    ".sql", ".css", ".scss", ".less",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".xml", ".svg",
}
CSV_EXTENSIONS = {".csv"}
JSON_EXTENSIONS = {".json"}
HTML_EXTENSIONS = {".html", ".htm"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp"}
DOCX_EXTENSION = {".docx"}
PDF_EXTENSION = {".pdf"}
PPTX_EXTENSION = {".pptx"}
XLSX_EXTENSIONS = {".xlsx", ".xls"}

ALL_SUPPORTED = (
    TEXT_EXTENSIONS | CSV_EXTENSIONS | JSON_EXTENSIONS
    | HTML_EXTENSIONS | IMAGE_EXTENSIONS
    | DOCX_EXTENSION | PDF_EXTENSION
    | PPTX_EXTENSION | XLSX_EXTENSIONS
)

# ===== 工具函数 =====
def _require_fitz():
    if fitz is None:
        raise RuntimeError("缺少依赖 PyMuPDF，请先安装。")

def _require_docx():
    if Document is None:
        raise RuntimeError("缺少依赖 python-docx，请先安装。")

def _require_pptx():
    if Presentation is None:
        raise RuntimeError("缺少依赖 python-pptx，请先安装。")

def _require_openpyxl():
    if load_workbook is None:
        raise RuntimeError("缺少依赖 openpyxl，请先安装。")

def _require_text_splitter():
    if RecursiveCharacterTextSplitter is None:
        raise RuntimeError("缺少依赖 langchain-text-splitters，请先安装。")


# ===== TXT / 纯文本 / 源代码 =====
def load_text(path: str | Path) -> str:
    """加载纯文本文件（含源代码、Markdown等）。"""
    return Path(path).read_text(encoding="utf-8", errors="ignore")


# ===== CSV =====
def load_csv(path: str | Path) -> str:
    """加载 CSV 文件，转成自然语言表格描述。"""
    lines = []
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        headers = None
        for i, row in enumerate(reader):
            if i == 0:
                headers = row
                lines.append(" | ".join(headers))
            else:
                if headers and row:
                    pairs = [f"{h}: {v}" for h, v in zip(headers, row) if v.strip()]
                    lines.append(" | ".join(pairs))
                elif row:
                    lines.append(" | ".join(row))
    return "\n".join(lines)


# ===== JSON =====
def load_json(path: str | Path) -> str:
    """加载 JSON 文件，转成格式化文本。"""
    data = json.loads(Path(path).read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(data, ensure_ascii=False, indent=2)


# ===== HTML =====
class _HTMLStripper(html.parser.HTMLParser):
    def __init__(self):
        super().__init__()
        self.text = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style", "noscript"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style", "noscript"):
            self._skip = False
        if tag in ("p", "br", "div", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
            self.text.append("\n")

    def handle_data(self, data):
        if not self._skip:
            stripped = data.strip()
            if stripped:
                self.text.append(stripped + " ")

    def get_text(self):
        return re.sub(r"\n{3,}", "\n\n", "".join(self.text)).strip()


def load_html(path: str | Path) -> str:
    """加载 HTML，剥离标签保留纯文本。"""
    content = Path(path).read_text(encoding="utf-8", errors="ignore")
    stripper = _HTMLStripper()
    stripper.feed(content)
    return stripper.get_text()


# ===== DOCX =====
def _iter_block_items(parent: DocumentType | Table):
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    if isinstance(parent, DocumentType):
        parent_elm = parent.element.body
    else:
        parent_elm = parent._tc

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _extract_images_from_docx(doc: DocumentType) -> list[bytes]:
    images: list[bytes] = []
    part = doc.part
    seen: set[str] = set()
    for rel in part.rels.values():
        reltype = getattr(rel, "reltype", "")
        if "image" not in reltype:
            continue
        image_part = getattr(rel, "target_part", None)
        if image_part is None:
            continue
        partname = str(getattr(image_part, "partname", ""))
        if partname in seen:
            continue
        seen.add(partname)
        blob = getattr(image_part, "blob", b"")
        if blob:
            images.append(blob)
    return images


def load_docx(path: str | Path) -> tuple[str, list[bytes]]:
    _require_docx()
    doc = Document(str(path))
    parts: list[str] = []

    for block in _iter_block_items(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if text:
                parts.append(text)
        elif isinstance(block, Table):
            for row in block.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)

    return "\n".join(parts), _extract_images_from_docx(doc)


# ===== PDF =====
def _extract_images_from_pdf(doc: Any, page_num: int) -> list[bytes]:
    page = doc.load_page(page_num)
    images: list[bytes] = []
    for image_info in page.get_images(full=True):
        xref = image_info[0]
        extracted = doc.extract_image(xref)
        image_bytes = extracted.get("image", b"")
        if image_bytes:
            images.append(image_bytes)
    return images


def load_pdf(path: str | Path) -> tuple[str, list[bytes]]:
    _require_fitz()
    doc = fitz.open(str(path))
    pages: list[str] = []
    images: list[bytes] = []
    try:
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            text = page.get_text().strip()
            if text:
                pages.append(f"[第 {page_num + 1} 页]\n{text}")
            images.extend(_extract_images_from_pdf(doc, page_num))
    finally:
        doc.close()
    return "\n\n".join(pages), images


# ===== PPTX =====
def load_pptx(path: str | Path) -> str:
    _require_pptx()
    prs = Presentation(str(path))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"[幻灯片 {slide_num}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_parts.append(row_text)
        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


# ===== Excel =====
def load_excel(path: str | Path) -> str:
    _require_openpyxl()
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts: list[str] = [f"[工作表: {sheet_name}]"]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                sheet_parts.append(row_text)
        parts.append("\n".join(sheet_parts))

    wb.close()
    return "\n\n".join(parts)


# ===== 图片 =====
def _describe_images(images: list[bytes]) -> list[str]:
    descriptions: list[str] = []
    for idx, image in enumerate(images, start=1):
        description = describe_image(image)
        descriptions.append(f"[图片{idx}内容]\n{description}")
    return descriptions


# ===== 切分 =====
def _split_text(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    _require_text_splitter()
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", "。", "！", "？", "；", "，", " ", ""],
    )
    return [chunk.strip() for chunk in splitter.split_text(text) if chunk.strip()]


# ===== 主入口 =====
def process_document(file_path: str | Path, chunk_size: int, chunk_overlap: int) -> list[str]:
    path = Path(file_path)
    suffix = path.suffix.lower()

    # 纯文本 / 源代码
    if suffix in TEXT_EXTENSIONS:
        text = load_text(path)
        return _split_text(text, chunk_size, chunk_overlap)

    # CSV
    if suffix in CSV_EXTENSIONS:
        text = load_csv(path)
        return _split_text(text, chunk_size, chunk_overlap)

    # JSON
    if suffix in JSON_EXTENSIONS:
        text = load_json(path)
        return _split_text(text, chunk_size, chunk_overlap)

    # HTML
    if suffix in HTML_EXTENSIONS:
        text = load_html(path)
        return _split_text(text, chunk_size, chunk_overlap)

    # DOCX（支持图片提取 → 多模态识别）
    if suffix in DOCX_EXTENSION:
        text, images = load_docx(path)
        image_text = "\n\n".join(_describe_images(images))
        combined = "\n\n".join(part for part in [text, image_text] if part.strip())
        return _split_text(combined, chunk_size, chunk_overlap) if combined.strip() else []

    # PDF（支持图片提取 → 多模态识别）
    if suffix in PDF_EXTENSION:
        text, images = load_pdf(path)
        image_text = "\n\n".join(_describe_images(images))
        combined = "\n\n".join(part for part in [text, image_text] if part.strip())
        return _split_text(combined, chunk_size, chunk_overlap) if combined.strip() else []

    # PPTX
    if suffix in PPTX_EXTENSION:
        text = load_pptx(path)
        return _split_text(text, chunk_size, chunk_overlap)

    # Excel
    if suffix in XLSX_EXTENSIONS:
        text = load_excel(path)
        return _split_text(text, chunk_size, chunk_overlap)

    # 图片（单张图片直接多模态识别）
    if suffix in IMAGE_EXTENSIONS:
        image_bytes = path.read_bytes()
        return [describe_image(image_bytes)]

    raise ValueError(f"暂不支持的文件类型: {suffix}")
